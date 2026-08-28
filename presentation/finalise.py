# -*- coding: utf-8 -*-
"""
Post-build step: stamp each speaker note with its time budget and the running
clock, so the deck can be rehearsed against the 10-minute slot without a
separate crib sheet. Also reports any slide whose script does not fit.
"""
import sys
from pptx import Presentation
from pptx.util import Pt

PPTX = sys.argv[1]
WPM = 140

# seconds budgeted per content slide (slides 1..12); the backup section is 0
BUDGET = [20, 43, 47, 72, 49, 57, 75, 63, 54, 39, 46, 33]

prs = Presentation(PPTX)
elapsed = 0
problems = []

for i, slide in enumerate(prs.slides):
    if not slide.has_notes_slide:
        continue
    tf = slide.notes_slide.notes_text_frame
    text = tf.text
    if not text.strip():
        continue
    words = len(text.split())
    spoken = words / WPM * 60

    if i < len(BUDGET):
        budget = BUDGET[i]
        elapsed += budget
        mm, ss = divmod(elapsed, 60)
        head = (f"[slide {i+1}  ·  {budget}s  ·  ends {mm}:{ss:02d} of 10:00  "
                f"·  {words} words ≈ {spoken:.0f}s at {WPM} wpm]")
        if spoken > budget * 1.10:
            problems.append(f"slide {i+1}: {spoken:.0f}s of script in a {budget}s slot")
        elif spoken < budget * 0.70:
            problems.append(f"slide {i+1}: only {spoken:.0f}s of script in a {budget}s slot")
    else:
        head = "[backup slide — not part of the 10 minutes]"

    tf.text = head + "\n\n" + text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)

prs.save(PPTX)

# python-pptx's round-trip renames the embedded media but does not always carry
# the matching <Default Extension=...> across, which leaves PowerPoint with an
# image part it has no content type for. Put back any that went missing.
import zipfile, shutil, os, re

CT = "[Content_Types].xml"
TYPES = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
         "gif": "image/gif", "emf": "image/x-emf", "wmf": "image/x-wmf"}

with zipfile.ZipFile(PPTX) as z:
    names = z.namelist()
    ct = z.read(CT).decode("utf-8")
    parts = {n: z.read(n) for n in names}

needed = {os.path.splitext(n)[1][1:].lower() for n in names if n.startswith("ppt/media/")}
missing = [e for e in sorted(needed)
           if e in TYPES and f'Extension="{e}"' not in ct]

if missing:
    add = "".join(f'<Default Extension="{e}" ContentType="{TYPES[e]}"/>' for e in missing)
    ct = re.sub(r"(<Types[^>]*>)", r"\1" + add, ct, count=1)
    parts[CT] = ct.encode("utf-8")
    tmp = PPTX + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, parts[n])
    shutil.move(tmp, PPTX)
    print(f"  repaired content types: {', '.join(missing)}")

print(f"stamped {PPTX}: {elapsed}s budgeted across {len(BUDGET)} content slides")
for p in problems:
    print("  ! " + p)
if not problems:
    print("  every script fits its slot")
